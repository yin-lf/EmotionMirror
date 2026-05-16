# -*- coding: utf-8 -*-
"""Update Slide 4 and Slide 6 content in Emotion汇报.pptx"""
from pptx import Presentation
from pptx.util import Pt

prs = Presentation('/home/ylf/EmotionMirror/Emotion汇报.pptx')

def get_textbox_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name and shape.has_text_frame:
            return shape
    return None

def set_textbox_content(shape, paragraphs):
    """Replace all paragraphs. paragraphs = [(text, bold), ...]"""
    tf = shape.text_frame
    for i, (text, bold) in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(4)
    xml_el = tf._txBody
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    all_p = xml_el.findall('a:p', nsmap)
    for extra in all_p[len(paragraphs):]:
        xml_el.remove(extra)

# ═══════════════════════════════════════
# Slide 4: 研究背景与现实意义
# ═══════════════════════════════════════
slide4 = prs.slides[3]

shape = get_textbox_by_name(slide4, '文本框 2')
if shape:
    set_textbox_content(shape, [
        ("Motivation：", True),
        ("随着人机交互技术的快速发展，计算机不再仅仅是工具，更逐步承担起“数字伙伴”的角色。"
         "然而，当前大多数人机交互系统缺乏对用户情绪状态的感知能力，导致交互过程机械、冰冷，难以建立情感层面的连接。"
         "随着虚拟数字人、AIGC 与智能陪伴系统的发展，用户对于“具有情感反馈能力”的 AI 需求快速提升。"
         "本作品旨在设计并实现一个基于多模态情绪感知的数字分身桌宠系统——Emotion Mirror，"
         "通过文本、语音两种模态实时感知用户情绪，驱动虚拟数字分身进行情绪化的动态反馈交互。", False),
    ])

shape = get_textbox_by_name(slide4, '文本框 40')
if shape:
    set_textbox_content(shape, [
        ("This Work Proposes a EmotionMirror：", True),
        ("这是一个开箱即用的端到端多模态情绪感知数字分身系统。"
         "用户通过文本或语音输入情绪内容，系统自动完成情绪识别、情绪向量生成、个性化表情驱动全流程，"
         "最终输出带有3D微动效果的数字人情绪镜像，并且可以在桌面进行操作。"
         "系统还集成了基于 DeepSeek 大语言模型的智能对话能力，桌宠可根据用户情绪和对话内容生成个性化回复，"
         "实现情感感知与智能陪伴的有机结合。"
         "整个链路无需手动调参，上传一张照片即可获得千人千面的情绪化数字分身。", False),
    ])

# ═══════════════════════════════════════
# Slide 6: 项目创新性
# ═══════════════════════════════════════
slide6 = prs.slides[5]

shape = get_textbox_by_name(slide6, '文本框 2')
if shape:
    set_textbox_content(shape, [
        ("模态统一情绪向量表示：", True),
        ("支持文本和语音两种异构输入模态，分别通过规则驱动引擎和 LSTM 模型进行情绪识别，"
         "两种模态的输出统一映射为 V-A-D 三维情绪向量 + 8 类中文情绪标签，"
         "在同一向量空间中实现标准化情绪表示，为后续扩展更多模态提供即插即用的接口规范。", False),
        ("个性化数字分身定制：", True),
        ("用户可上传任意个人照片（自拍/二次元/卡通），系统通过 rembg 自动去除背景、"
         "LivePortrait 提取人脸特征并重定向表情，实现千人千面的情绪化数字分身，突破了预设表情素材库的局限性。"
         "支持情绪强度 1-5 级可调，以及微笑、眉毛、唇部变化等 9 个底层参数的精细微调。", False),
        ("端到端的数字人情绪表示系统：", True),
        ("区别于传统情感分析系统仅输出分类标签的做法，本作品将情绪识别结果直接映射为 LivePortrait 面部动作参数"
         "（微笑、眉毛、眼球方向、唇部变化等），基于用户上传的个人照片生成个性化的表情动画 GIF，"
         "并叠加余弦缓动插值与头部微摆效果，输出自然流畅的透明背景动画。", False),
        ("便携的桌宠交互系统：", True),
        ("基于 PySide6 实现桌面宠物窗口，支持情绪 GIF 展示、十字淡入淡出动画切换、"
         "以及基于 DeepSeek 大语言模型的智能对话功能。"
         "桌宠可根据用户的文字输入实时分析情绪，同步更新表情动画，并生成带有情绪标签的个性化回复，"
         "形成“输入-识别-表情驱动-智能回复”的完整情感交互闭环。", False),
        ("沉浸式场景渲染：", True),
        ("前端根据系统时间（早晨/下午/傍晚/夜晚四时段）与当前情绪自动匹配场景背景，"
         "共 4×8=32 种渐变场景组合，搭配雨滴、星星、火花、雾气、闪电等 6 种 CSS 动画效果，"
         "为表情展示提供沉浸式氛围杖托。", False),
    ])

out = '/home/ylf/EmotionMirror/Emotion汇报_new.pptx'
prs.save(out)
print(f'Saved to {out}')
