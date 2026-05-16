"""Generate PPT for EmotionMirror expression synthesis submodule."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD   = RGBColor(0x16, 0x21, 0x3E)
ACCENT    = RGBColor(0x25, 0x63, 0xEB)
ACCENT2   = RGBColor(0x38, 0xBD, 0xF8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xA0, 0xAE, 0xC0)
ORANGE    = RGBColor(0xF0, 0x98, 0x19)
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
RED       = RGBColor(0xE5, 0x3E, 0x3E)
PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(6)


def add_table(slide, left, top, width, height, rows, cols, data, header_color=ACCENT):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Microsoft YaHei'
                paragraph.alignment = PP_ALIGN.CENTER
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = WHITE

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else RGBColor(0x0F, 0x18, 0x2A)
    return table


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ═══════════════════════════════════════════
# Slide 1: Title
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             'EmotionMirror', font_size=48, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3), Inches(11), Inches(1),
             '情感合成子模块', font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
             '多模态情感驱动的人物表情合成与场景渲染', font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
             'backend/expression/', font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 2: Overview / Pipeline
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
             '整体流水线', font_size=32, color=ACCENT2, bold=True)

# Pipeline boxes
pipeline_steps = [
    ('文本/语音\n情感识别', ACCENT, Inches(0.5)),
    ('VAD 情感向量\n[V, A, D]', RGBColor(0x8B, 0x5C, 0xF6), Inches(2.8)),
    ('情绪标签\n→ 表情参数', ORANGE, Inches(5.1)),
    ('LivePortrait\n表情重定向', GREEN, Inches(7.4)),
    ('GIF 生成\n透明背景', RED, Inches(9.7)),
    ('场景背景\n时间×情绪', RGBColor(0x06, 0xB6, 0xD4), Inches(12.0)),
]

for label, color, left in pipeline_steps:
    shape = add_card(slide, left, Inches(1.8), Inches(2.0), Inches(1.2), color)
    shape.text_frame.word_wrap = True
    for p in shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    tf = shape.text_frame
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Arrows between boxes
for i in range(len(pipeline_steps) - 1):
    _, _, left1 = pipeline_steps[i]
    _, _, left2 = pipeline_steps[i + 1]
    arrow_left = left1 + Inches(2.0)
    arrow_width = left2 - arrow_left
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, int(arrow_left), Inches(2.15), int(arrow_width), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()

# Bottom description
desc_items = [
    '• 文本分析：基于关键词规则 + VAD 向量映射（8种情绪，95个关键词）',
    '• 语音分析：librosa 312维特征 → LSTM(128) → softmax 6分类 → VAD 映射',
    '• 表情合成：EMOTION_PARAMS 查表 → 强度缩放 → LivePortrait 关键点驱动 → SPADE 解码',
    '• GIF 输出：24帧余弦插值 + 头部微摆 + rembg 透明化 + 255色调色板',
    '• 场景渲染：4时段 × 8情绪 = 32种渐变背景 + 6种 CSS 动画效果',
]
add_bullet_list(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.5),
                desc_items, font_size=17, color=WHITE)

# ═══════════════════════════════════════════
# Slide 3: VAD Vector Mapping
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
             'VAD 情感向量 → 面部动作参数映射', font_size=32, color=ACCENT2, bold=True)

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             'VAD 向量 = [效价 Valence, 唤醒度 Arousal, 优势度 Dominance]，归一化至 [0, 1]',
             font_size=15, color=GRAY)

# VAD table
vad_data = [
    ['情绪', 'V (效价)', 'A (唤醒度)', 'D (优势度)', '对应表情参数'],
    ['开心', '0.86', '0.72', '0.62', 'smile=1.3, lip_two=12, eyebrow=5'],
    ['悲伤', '0.18', '0.32', '0.38', 'smile=-0.3, eyebrow=-25, lip_three=-60'],
    ['愤怒', '0.12', '0.78', '0.67', 'smile=-0.3, eyebrow=-30, lip_three=-70'],
    ['焦虑', '0.25', '0.82', '0.35', 'eyebrow=-20, lip_one=-12, lip_three=-30'],
    ['恐惧', '0.18', '0.86', '0.25', 'eyebrow=30, lip_one=15, lip_three=80'],
    ['平静', '0.52', '0.22', '0.55', '(全部默认值 0)'],
    ['厌恶', '0.10', '0.55', '0.42', 'smile=-0.3, lip_zero=-0.08, lip_three=-40'],
    ['惊讶', '0.63', '0.74', '0.48', 'smile=0.5, eyebrow=30, lip_three=100'],
]
add_table(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.8), 9, 5, vad_data)

add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
             'synthesis.py: EMOTION_PARAMS — 情绪标签直接查表映射为 LivePortrait 的 9 个面部动作参数',
             font_size=13, color=ACCENT)

# ═══════════════════════════════════════════
# Slide 4: LivePortrait Architecture
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
             '基于 LivePortrait 的人物表情合成', font_size=32, color=ACCENT2, bold=True)

# Left: Architecture
add_card(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.5), BG_CARD)
add_text_box(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.5),
             'LivePortrait 重定向流程', font_size=20, color=ORANGE, bold=True)

arch_items = [
    '1. init_retargeting_image()',
    '   → 人脸检测 + 裁剪 (scale=2.3)',
    '   → 提取外观特征 f_s + 关键点 x_s',
    '   → 计算旋转矩阵 R_s, R_d',
    '',
    '2. execute_image_retargeting()',
    '   → 对 21 个面部关键点应用 delta 参数:',
    '     • update_delta_new_smile  → 嘴角/脸颊',
    '     • update_delta_new_eyebrow → 眉毛',
    '     • update_delta_new_lip_*   → 嘴唇',
    '     • update_delta_new_eyeball → 眼球/眼睑',
    '',
    '3. 关键点重建:',
    '   x_d = scale × (x_c @ R_d + delta) + t',
    '',
    '4. stitching() → warp_decode() → paste_back()',
    '   → SPADE 生成器 + 扭曲网络 → 回贴原图',
]
add_bullet_list(slide, Inches(0.7), Inches(2.2), Inches(5.6), Inches(4.6),
                arch_items, font_size=14, color=WHITE)

# Right: Parameter ranges
add_card(slide, Inches(6.9), Inches(1.5), Inches(6.0), Inches(5.5), BG_CARD)
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.6), Inches(0.5),
             '表情参数范围 (PARAM_RANGES)', font_size=20, color=ORANGE, bold=True)

param_data = [
    ['参数', '范围', '步长', '说明'],
    ['smile', '-1.0 ~ 1.5', '0.1', '微笑/嘴角上扬'],
    ['wink', '0 ~ 1.0', '0.05', '眨眼'],
    ['eyebrow', '-40 ~ 40', '1', '眉毛升降'],
    ['eyeball_x', '-60 ~ 60', '1', '眼球水平方向'],
    ['eyeball_y', '-60 ~ 60', '1', '眼球垂直方向'],
    ['lip_zero', '-0.1 ~ 0.1', '0.005', '嘴唇微调'],
    ['lip_one', '-30 ~ 30', '1', '嘴唇张合'],
    ['lip_two', '-30 ~ 30', '1', '嘴角变化'],
    ['lip_three', '-100 ~ 100', '1', '嘴唇大幅变化'],
]
add_table(slide, Inches(7.0), Inches(2.3), Inches(5.7), Inches(4.5), 10, 4, param_data, ORANGE)

# ═══════════════════════════════════════════
# Slide 5: Emotion Intensity Control
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
             '情绪强度控制', font_size=32, color=ACCENT2, bold=True)

# Left: Intensity mechanism
add_card(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.5), BG_CARD)
add_text_box(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.5),
             '_apply_intensity() 参数缩放', font_size=20, color=GREEN, bold=True)

intensity_items = [
    '• 用户可调节情绪强度：1 ~ 5 级（默认5=满强度）',
    '• 缩放公式：scale = intensity / 5.0',
    '  所有表情参数 × scale → 等比缩放',
    '',
    '  intensity=1 → 20%    intensity=2 → 40%',
    '  intensity=3 → 60%    intensity=4 → 80%',
    '  intensity=5 → 100% (原始值)',
    '',
    '• 合并策略：user_params 覆盖默认值',
    '  merged = {**EMOTION_PARAMS[emotion], **user_params}',
    '  merged = _apply_intensity(merged, intensity)',
    '',
    '• 前端 UI：简洁强度滑块 + 可折叠高级面板',
    '  高级面板展开后可分别调节 9 个底层参数',
]
add_bullet_list(slide, Inches(0.7), Inches(2.2), Inches(5.6), Inches(4.6),
                intensity_items, font_size=15, color=WHITE)

# Right: Code snippet
add_card(slide, Inches(6.9), Inches(1.5), Inches(6.0), Inches(5.5), RGBColor(0x0D, 0x12, 0x1A))
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.6), Inches(0.5),
             '核心代码', font_size=20, color=GREEN, bold=True)

code_items = [
    'def _apply_intensity(base, intensity):',
    '    if not base or intensity >= 5:',
    '        return base',
    '    s = max(1, min(5, intensity)) / 5.0',
    '    return {k: v * s for k, v in base.items()}',
    '',
    'def synthesize_expression_gif(',
    '    image_path, emotion,',
    '    num_frames=12, fps=10,',
    '    params=None, intensity=5  # ← 新增',
    '):',
    '    base = EMOTION_PARAMS.get(emotion, {})',
    '    base = _apply_intensity(base, intensity)',
    '    target = {**base, **params} if params else base',
    '    ...',
]
add_bullet_list(slide, Inches(7.1), Inches(2.2), Inches(5.6), Inches(4.6),
                code_items, font_size=14, color=ACCENT2)

# ═══════════════════════════════════════════
# Slide 6: GIF Generation
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
             'GIF 动画生成与头部微摆', font_size=32, color=ACCENT2, bold=True)

# Left: Frame interpolation
add_card(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.5), BG_CARD)
add_text_box(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.5),
             '帧插值策略', font_size=20, color=PURPLE, bold=True)

gif_items = [
    '• 总帧数：2 × num_frames = 24 帧',
    '  Phase 1: 中性 → 表情 (12帧, 渐入)',
    '  Phase 2: 表情 → 中性 (12帧, 渐出)',
    '',
    '• 余弦缓动函数 (cosine ease-in-out):',
    '  t = 0.5 - 0.5 × cos(t × π)',
    '  使起止帧平滑过渡，无突变',
    '',
    '• 无缝循环：末帧 t→0 与首帧 t=0 一致',
    '',
    '• 输出：10fps, 2.4秒循环, loop=0(无限)',
    '  透明背景: rembg抠图 + 255色调色板',
    '  disposal=2 确保帧间正确清除',
]
add_bullet_list(slide, Inches(0.7), Inches(2.2), Inches(5.6), Inches(4.6),
                gif_items, font_size=15, color=WHITE)

# Right: Head pose
add_card(slide, Inches(6.9), Inches(1.5), Inches(6.0), Inches(5.5), BG_CARD)
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.6), Inches(0.5),
             '头部微摆 (Pseudo-3D)', font_size=20, color=PURPLE, bold=True)

pose_items = [
    '为静态表情添加自然的头部晃动：',
    '',
    'phase = 2π × frame_idx / total_frames',
    'yaw   = 1.2 × sin(1 × phase)',
    'pitch = 0.8 × sin(phase + π/4)',
    'roll  = 0.5 × sin(1 × phase)',
    '',
    '• 整数倍频率保证首尾帧一致 (无缝循环)',
    '• pitch 偏移 π/4 避免机械感',
    '• 幅度经调优：yaw=1.2°, pitch=0.8°, roll=0.5°',
    '',
    '透明背景处理：',
    '• 首帧通过 rembg 生成 alpha 蒙版',
    '• 高度≥200px 时形态学膨胀防锯齿',
    '• 非前景像素 alpha=0',
]
add_bullet_list(slide, Inches(7.1), Inches(2.2), Inches(5.6), Inches(4.6),
                pose_items, font_size=15, color=WHITE)

# ═══════════════════════════════════════════
# Slide 7: Scene Background
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
             '基于时间与情绪的场景背景', font_size=32, color=ACCENT2, bold=True)

add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             '4时段 × 8情绪 = 32种场景组合，每种包含 CSS 渐变背景 + 动画效果',
             font_size=16, color=GRAY)

# Time period cards
periods = [
    ('早晨 (6-12)', '暖橙/粉色系', RGBColor(0xFF, 0x9A, 0x56)),
    ('下午 (12-18)', '明亮蓝天系', RGBColor(0x74, 0xB9, 0xFF)),
    ('傍晚 (18-21)', '夕阳/暗橙系', RGBColor(0xF0, 0x93, 0xFB)),
    ('夜晚 (21-6)', '深蓝/深紫系', RGBColor(0x0F, 0x20, 0x27)),
]
for i, (title, desc, color) in enumerate(periods):
    left = Inches(0.5 + i * 3.2)
    add_card(slide, left, Inches(1.9), Inches(2.9), Inches(1.0), color)
    tf = slide.shapes[-1].text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE if i == 3 else BG_DARK
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = WHITE if i == 3 else BG_DARK
    p2.font.name = 'Microsoft YaHei'
    p2.alignment = PP_ALIGN.CENTER

# Animation types
add_card(slide, Inches(0.4), Inches(3.2), Inches(12.5), Inches(4.0), BG_CARD)
add_text_box(slide, Inches(0.6), Inches(3.3), Inches(12), Inches(0.5),
             '6种动画效果', font_size=20, color=ACCENT2, bold=True)

anim_data = [
    ['动画', '实现方式', '适用场景', 'CSS 技术'],
    ['rain (雨滴)', '12个半透明细条从上到下', '悲伤情绪 (全时段)', '@keyframes + 线性动画'],
    ['stars (星星)', '8个圆点闪烁 (0.2↔1.0 透明度)', '夜晚 开心/平静', 'ease-in-out 循环'],
    ['fire (火花)', '10个渐变点从下往上', '愤怒情绪 (全时段)', 'ease-out + 缩放'],
    ['fog (雾气)', '3个大模糊圆水平漂移', '焦虑情绪 (全时段)', 'blur(20px) + alternate'],
    ['lightning (闪电)', '全屏白色闪光 (4秒周期)', '恐惧情绪 (全时段)', '关键帧精确控制'],
    ['sparkle (光斑)', '6个光点缩放闪烁', '开心/惊讶 (全时段)', 'ease-in-out + scale'],
]
add_table(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(3.0), 7, 4, anim_data, ACCENT)

# ═══════════════════════════════════════════
# Slide 8: Demo / Summary
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
             '总结', font_size=32, color=ACCENT2, bold=True)

summary_items = [
    '1. 完整的情感→表情流水线',
    '   文本/语音情感识别 → VAD向量 → EMOTION_PARAMS查表 → LivePortrait关键点驱动 → GIF输出',
    '',
    '2. 灵活的情绪强度控制',
    '   1-5级等比缩放 + 高级参数面板(9个独立滑块) → 前后端完整联调',
    '',
    '3. 自然的动画效果',
    '   余弦缓动插值 + 头部微摆(yaw/pitch/roll) + rembg透明背景',
    '',
    '4. 沉浸式场景渲染',
    '   4时段 × 8情绪 = 32种场景，6种CSS动画(雨/星/火/雾/电/光)',
    '',
    '5. 可扩展设计',
    '   EMOTION_PARAMS / PARAM_RANGES / SCENES 配置驱动，新增情绪或场景只需修改配置表',
]
add_bullet_list(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(5.5),
                summary_items, font_size=18, color=WHITE)

# Key files
add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
             '核心文件: synthesis.py | sceneConfig.js | SceneBackground.jsx | ExpressionParamsPanel.jsx | api.js',
             font_size=13, color=ACCENT)

# ═══════════════════════════════════════════
# Save
# ═══════════════════════════════════════════
out = '/home/ylf/EmotionMirror/情感合成子模块.pptx'
prs.save(out)
print(f'Saved to {out}')
